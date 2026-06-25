"""Assemble Pixxel board deck PDF and PPTX from slide PNGs in deliverables/slides/."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

ROOT = Path(__file__).resolve().parents[1]
DELIVERABLES = ROOT / "deliverables"
SLIDES_DIR = DELIVERABLES / "slides"

SLIDE_NAMES = [
    "01_executive_summary.png",
    "02_tam_sizing.png",
    "03_segment_selection.png",
    "04_monetization_model.png",
    "05_roadmap_next_steps.png",
]


def collect_slides() -> list[Path]:
    slides: list[Path] = []
    for name in SLIDE_NAMES:
        path = SLIDES_DIR / name
        if not path.exists():
            raise FileNotFoundError(f"Missing slide: {path}")
        slides.append(path)
    return slides


def build_pdf(slide_paths: list[Path], pdf_path: Path) -> None:
    first = Image.open(slide_paths[0])
    page_w, page_h = first.size
    first.close()

    c = canvas.Canvas(str(pdf_path), pagesize=(page_w, page_h))
    for slide in slide_paths:
        c.drawImage(ImageReader(str(slide)), 0, 0, width=page_w, height=page_h)
        c.showPage()
    c.save()


def build_pptx(slide_paths: list[Path], pptx_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for slide_path in slide_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(slide_path), 0, 0, width=prs.slide_width, height=prs.slide_height
        )

    prs.save(str(pptx_path))


def main() -> None:
    DELIVERABLES.mkdir(parents=True, exist_ok=True)
    slides = collect_slides()

    pdf_path = DELIVERABLES / "Pixxel_Commercial_Strategy_Deck.pdf"
    pptx_path = DELIVERABLES / "Pixxel_Commercial_Strategy_Deck.pptx"

    build_pdf(slides, pdf_path)
    build_pptx(slides, pptx_path)

    print(f"Slides: {SLIDES_DIR}")
    for s in slides:
        print(f"  - {s.name}")
    print(f"PDF:  {pdf_path}")
    print(f"PPTX: {pptx_path}")


if __name__ == "__main__":
    main()
